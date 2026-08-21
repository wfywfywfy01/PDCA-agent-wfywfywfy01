# -*- coding: utf-8 -*-
from pptx import Presentation

out_path = r"d:\经销商PDCA\_tmp_pptx_content.txt"
path = r"c:\Users\frank\Desktop\海外经销商2026年7月第二周周报.pptx"
prs = Presentation(path)

lines = []
lines.append(f"Slides: {len(prs.slides)}")
lines.append(f"Slide size: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches")

for i, slide in enumerate(prs.slides):
    lines.append(f"\n===== SLIDE {i+1} =====")
    for shape in slide.shapes:
        lines.append(f"  shape: {shape.shape_type}, name={shape.name}")
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = p.text.strip()
                if t:
                    lines.append(f"    TEXT: {t}")
        if shape.has_table:
            table = shape.table
            lines.append(f"    TABLE {len(table.rows)}x{len(table.columns)}:")
            for r_idx, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("\n", " | ") for cell in row.cells]
                lines.append(f"      R{r_idx}: {' | '.join(cells)}")

with open(out_path, "w", encoding="utf-8") as f:
    f.write("\n".join(lines))
print(f"Wrote {len(lines)} lines to {out_path}")
