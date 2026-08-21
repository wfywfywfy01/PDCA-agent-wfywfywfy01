# -*- coding: utf-8 -*-
"""Extract text/tables from team W3 PPTX files into one markdown dump."""
from __future__ import annotations

import sys
from pathlib import Path

sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation

FILES = [
    Path(r"c:\Users\frank\xwechat_files\wxid_ewjxd3lystvp22_2fcf\msg\file\2026-07\周报(1).pptx"),
    Path(r"c:\Users\frank\xwechat_files\wxid_ewjxd3lystvp22_2fcf\msg\file\2026-07\Lina2026年7月第三周周报(1).pptx"),
    Path(r"c:\Users\frank\xwechat_files\wxid_ewjxd3lystvp22_2fcf\msg\file\2026-07\第三周周报7月(1).pptx"),
    Path(r"c:\Users\frank\xwechat_files\wxid_ewjxd3lystvp22_2fcf\msg\file\2026-07\刘雪梅7月第三周周报(1).pptx"),
]
OUT = Path(r"d:\经销商PDCA\overseas_weekly\outputs\_w3_team_ppts_extract.md")


def extract(path: Path) -> list[str]:
    lines = [f"# FILE: {path.name}", f"path: {path}", ""]
    if not path.exists():
        lines.append("**MISSING**")
        return lines
    prs = Presentation(str(path))
    lines.append(f"Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        lines.append(f"\n## SLIDE {i+1}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join([r.text for r in p.runs]) if p.runs else p.text
                    if t.strip():
                        lines.append(t.strip())
            if shape.has_table:
                table = shape.table
                lines.append(f"[TABLE {len(table.rows)}x{len(table.columns)}]")
                for row in table.rows:
                    cells = [c.text.replace("\n", " / ").strip() for c in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
    return lines


def main() -> None:
    all_lines: list[str] = []
    for f in FILES:
        all_lines.extend(extract(f))
        all_lines.append("\n---\n")
    OUT.write_text("\n".join(all_lines), encoding="utf-8")
    print(OUT)
    print(f"bytes={OUT.stat().st_size} lines={len(all_lines)}")


if __name__ == "__main__":
    main()
